#!/usr/bin/env python3
"""
Script to generate PowerPoint presentation for IoT Anomaly Detection System
Run: pip install python-pptx && python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RgbColor(0, 51, 102)
LIGHT_BLUE = RgbColor(0, 112, 192)
ORANGE = RgbColor(237, 125, 49)
GREEN = RgbColor(84, 130, 53)
GRAY = RgbColor(89, 89, 89)
WHITE = RgbColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), prs.slide_width, Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, two_column=False, right_bullets=None):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if two_column and right_bullets:
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.space_after = Pt(12)
            p.level = 0

        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(right_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.space_after = Pt(12)
            p.level = 0
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY
            p.space_after = Pt(14)

    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Architecture boxes
    components = [
        ("IoT Sensors\n(Synthetic Data)", Inches(0.8), Inches(2), Inches(2), Inches(1), ORANGE),
        ("Apache Kafka\n(Message Broker)", Inches(3.5), Inches(2), Inches(2.2), Inches(1), LIGHT_BLUE),
        ("Detector Service\n(Isolation Forest)", Inches(6.5), Inches(2), Inches(2.2), Inches(1), GREEN),
        ("Kafka Topics\n(sensors.anomalies)", Inches(9.5), Inches(1.8), Inches(2.2), Inches(0.9), LIGHT_BLUE),
        ("HDFS Storage\n(Batch Data)", Inches(9.5), Inches(3), Inches(2.2), Inches(0.9), RgbColor(150, 75, 0)),
        ("Streamlit Dashboard\n(Visualization)", Inches(5), Inches(5), Inches(3), Inches(1.2), RgbColor(255, 75, 75)),
    ]

    for text, left, top, width, height, color in components:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow labels
    arrows = [
        ("sensors.raw", Inches(2.9), Inches(2.35), Inches(0.6)),
        ("Process", Inches(5.8), Inches(2.35), Inches(0.6)),
        ("Real-time", Inches(8.8), Inches(2.1), Inches(0.6)),
        ("Batch", Inches(8.8), Inches(3.3), Inches(0.6)),
    ]

    for text, left, top, width in arrows:
        box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"-> {text}"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY

    # Data flow description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.8))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Lambda Architecture: Real-time streaming (Kafka) + Batch processing (HDFS) for comprehensive data analysis"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_tech_table_slide(prs):
    """Add technology stack table slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table data
    table_data = [
        ("Layer", "Technology", "Purpose"),
        ("Data Source", "Python Generator", "Simulate IoT sensors"),
        ("Message Broker", "Apache Kafka", "Real-time stream ingestion"),
        ("Stream Processing", "Python + scikit-learn", "ML-based anomaly detection"),
        ("Batch Storage", "Apache Hadoop HDFS", "Historical data warehouse"),
        ("Visualization", "Streamlit + Plotly", "Interactive web dashboard"),
        ("Orchestration", "Docker Compose", "Container management"),
    ]

    # Create table
    rows = len(table_data)
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5)).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4)
    table.columns[2].width = Inches(4.833)

    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            if i == 0:  # Header row
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                p.font.color.rgb = GRAY
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RgbColor(240, 240, 240)

    return slide

def add_evaluation_slide(prs):
    """Add evaluation criteria mapping slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Evaluation Criteria Coverage"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column criteria
    left_items = [
        ("Data Producers (0.5p)", "Python synthetic IoT sensor generator"),
        ("Data Sources (0.5p)", "Temperature & Vibration sensors with anomaly injection"),
        ("Data Storage (1p)", "Apache Hadoop HDFS for batch storage"),
        ("Data Visualization (1p)", "Streamlit dashboard with Plotly charts"),
        ("Data Processing (1p)", "Real-time (Kafka) + Batch (HDFS) - Lambda Architecture"),
    ]

    right_items = [
        ("Solution Design (0.5p)", "Microservices with Docker Compose orchestration"),
        ("Prototype Maturity (0.5p)", "Fully automated, single-command deployment"),
        ("ML Algorithm (1p)", "Isolation Forest with adaptive retraining"),
        ("Project Demo (1p)", "Live dashboard with real-time anomaly detection"),
    ]

    # Left column
    y_pos = 1.5
    for title, desc in left_items:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(6), Inches(0.8))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        y_pos += 0.95

    # Right column
    y_pos = 1.5
    for title, desc in right_items:
        box = slide.shapes.add_textbox(Inches(6.8), Inches(y_pos), Inches(6), Inches(0.8))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        y_pos += 0.95

    return slide

# ============================================
# BUILD THE PRESENTATION
# ============================================

# Slide 1: Title
add_title_slide(prs,
    "IoT Anomaly Detection System",
    "Real-Time Streaming Pipeline for Industrial Sensor Monitoring\n\nTeam: Vlad, Mihai, Darius"
)

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement & Motivation", [
    "Challenge: Industrial facilities generate continuous streams of sensor data",
    "   - Equipment failures can be costly and dangerous",
    "   - Manual monitoring is impractical at scale",
    "   - Traditional threshold-based alerts miss complex patterns",
    "",
    "Our Solution: Automated anomaly detection using machine learning",
    "   - Real-time processing for immediate alerts",
    "   - Historical analysis for pattern discovery",
    "   - Scalable architecture for growing sensor networks",
    "",
    "Use Case: Manufacturing equipment monitoring",
    "   - Temperature sensors: Detect overheating machinery",
    "   - Vibration sensors: Identify mechanical wear/failures"
])

# Slide 3: Architecture Overview
add_architecture_slide(prs)

# Slide 4: Data Flow Pipeline
add_content_slide(prs, "Data Flow Pipeline", [
    "1. DATA GENERATION (Producer Service)",
    "   - Synthetic IoT sensors: Temperature (25C +/- 2C) and Vibration (5Hz +/- 2Hz)",
    "   - Anomaly injection: 5% probability of spikes (+/-15C or +20Hz)",
    "   - Throughput: ~20 events/second to Kafka",
    "",
    "2. STREAM INGESTION (Apache Kafka)",
    "   - Topic 'sensors.raw': Raw sensor readings",
    "   - Topic 'sensors.anomalies': Enriched data with ML predictions",
    "",
    "3. PROCESSING (Detector Service)",
    "   - Consumes from sensors.raw, enriches with features",
    "   - Publishes to sensors.anomalies (real-time path)",
    "   - Writes batches of 100 records to HDFS (batch path)",
    "",
    "4. VISUALIZATION (Streamlit Dashboard)",
    "   - Real-time monitoring with live charts",
    "   - Batch analysis from HDFS historical data"
])

# Slide 5: ML Algorithm
add_content_slide(prs, "Machine Learning: Isolation Forest", [
    "Algorithm Choice: Isolation Forest (Unsupervised Anomaly Detection)",
    "   - No labeled data required - learns normal patterns automatically",
    "   - Efficient for streaming: O(n log n) complexity",
    "   - Robust to outliers in training data",
    "",
    "Feature Engineering (6 features extracted per reading):",
    "   - rolling_mean: 10-point moving average (trend detection)",
    "   - rolling_std: 10-point standard deviation (volatility)",
    "   - velocity: First derivative (rate of change)",
    "   - acceleration: Second derivative (momentum)",
    "   - local_z: Z-score normalization (deviation strength)",
    "   - value: Raw sensor reading",
    "",
    "Adaptive Training:",
    "   - Separate model per sensor type (temperature, vibration)",
    "   - Retrains every 100 new data points",
    "   - Adapts to changing sensor behavior over time"
])

# Slide 6: Technology Stack
add_tech_table_slide(prs)

# Slide 7: Dashboard Features
add_content_slide(prs, "Dashboard Features",
    [
        "TAB 1: Real-Time Monitoring",
        "   - Live line charts for temperature & vibration",
        "   - Red markers overlay for anomalies",
        "   - Configurable refresh rate (0.5-5 sec)",
        "",
        "Advanced Feature Visualization:",
        "   - Bollinger Bands (trend + volatility)",
        "   - Velocity & acceleration derivatives",
        "   - Local z-score deviation charts",
    ],
    two_column=True,
    right_bullets=[
        "TAB 2: Batch Analysis (HDFS)",
        "   - Historical data from last 5 HDFS files",
        "   - Raw data table preview",
        "",
        "Statistical Analysis:",
        "   - Value distribution histograms",
        "   - Normal vs anomalous comparison",
        "   - Summary statistics by sensor type",
        "   - Feature correlation scatter matrix",
    ]
)

# Slide 8: Evaluation Criteria
add_evaluation_slide(prs)

# Slide 9: Team & Work Distribution
add_content_slide(prs, "Team Composition & Work Distribution", [
    "Team Members:",
    "",
    "   Vlad - Dashboard Development",
    "      - Streamlit web application",
    "      - Plotly visualizations",
    "      - HDFS data reader integration",
    "",
    "   Mihai - Data Producers",
    "      - Synthetic sensor simulation",
    "      - Kafka producer implementation",
    "      - Anomaly injection logic",
    "",
    "   Darius - Anomaly Detector (Predictor)",
    "      - Isolation Forest ML model",
    "      - Feature engineering pipeline",
    "      - Kafka consumer/producer integration",
    "",
    "   Joint Effort: Docker infrastructure, integration testing, documentation"
])

# Slide 10: Demo slide
add_content_slide(prs, "Live Demo", [
    "Starting the System:",
    "   $ docker-compose up --build -d",
    "",
    "Demo Flow:",
    "",
    "   1. Show Docker containers running (all 6 services)",
    "",
    "   2. Open Dashboard at http://localhost:8501",
    "",
    "   3. Real-Time Tab: Watch live sensor data streaming",
    "      - Observe normal readings and anomaly spikes (red markers)",
    "      - Show advanced feature charts (Bollinger bands, velocity)",
    "",
    "   4. Batch Analysis Tab: Load historical HDFS data",
    "      - View raw data table",
    "      - Show distribution histograms (normal vs anomaly)",
    "      - Display summary statistics"
])

# Slide 11: Thank you
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Background
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(32)
p.font.color.rgb = RgbColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# Team
team_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
tf = team_box.text_frame
p = tf.paragraphs[0]
p.text = "Vlad | Mihai | Darius"
p.font.size = Pt(20)
p.font.color.rgb = RgbColor(150, 150, 150)
p.alignment = PP_ALIGN.CENTER

# Save
prs.save('IoT_Anomaly_Detection_Presentation.pptx')
print("Presentation saved as 'IoT_Anomaly_Detection_Presentation.pptx'")
print(f"Total slides: {len(prs.slides)}")
